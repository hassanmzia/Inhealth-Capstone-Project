#!/usr/bin/env python3
"""
Generate InHealth Technical Architecture PowerPoint presentation.
Usage: python scripts/generate_architecture_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Helpers ──────────────────────────────────────────────────

def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False,
                color='#333333', alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None, text='',
             font_size=10, font_color='#FFFFFF', bold=False, radius=0.08):
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if border_color:
        shape.line.color.rgb = rgb(border_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for i, line in enumerate(text.split('\n')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = rgb(font_color)
            p.font.bold = bold if i == 0 else False
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
    return shape


# ── Main ─────────────────────────────────────────────────────

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ============================================================
    # SLIDE 1: Title
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_rect(slide, 0, 0, 13.333, 7.5, '0D47A1')
    add_textbox(slide, 1, 1.5, 11, 1.2, 'InHealth Chronic Care', 44, True, '#FFFFFF', PP_ALIGN.CENTER)
    add_textbox(slide, 1, 2.8, 11, 0.8,
                'Unified Autonomous Agentic AI Platform for Chronic Disease Management',
                22, False, '#BBDEFB', PP_ALIGN.CENTER)
    add_textbox(slide, 1, 4.0, 11, 0.5, 'Technical Architecture Overview', 18, False, '#90CAF9', PP_ALIGN.CENTER)
    add_textbox(slide, 1, 5.5, 11, 0.5,
                '430+ Files  |  20 Docker Services  |  25+ AI Agents  |  15 Django Apps  |  14 FHIR R4 Resources',
                14, False, '#64B5F6', PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 2: Architecture Overview
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 0.8, '0D47A1')
    add_textbox(slide, 0.3, 0.1, 12, 0.6, 'System Architecture Overview', 28, True, '#FFFFFF', PP_ALIGN.LEFT)

    # Layer boxes
    layers = [
        (0.3, 1.0, 12.7, 0.9, '4CAF50', 'PRESENTATION LAYER\nReact 18 + TypeScript + TailwindCSS + shadcn/ui + Recharts + Zustand + PWA'),
        (0.3, 2.0, 12.7, 0.5, 'FF9800', 'NGINX REVERSE PROXY (Port 8788)\n/ → Frontend  |  /api/ → Django  |  /ws/ → Channels  |  /agents/ → FastAPI  |  /mcp/ → MCP Server  |  /a2a/ → A2A Gateway'),
        (0.3, 2.6, 6.2, 1.6, '1976D2', 'BACKEND API — Django 5 + DRF + Channels\n15 Apps: accounts, tenants, fhir, hl7, patients, clinical,\nanalytics, notifications, mcp_bridge, a2a_bridge,\nresearch, telemedicine, billing, sdoh, dashboard\nML: LSTM, XGBoost, Random Forest, HMM, Digital Twin, Federated'),
        (6.8, 2.6, 6.2, 1.6, '1565C0', 'NODE.JS GATEWAY SERVICES\nMCP Server (3001): Tool Registry, Context Mgmt, LLM Proxy\nA2A Gateway (3002): Agent Cards, Task Delegation\nOpenTelemetry + Langfuse Tracing'),
        (0.3, 4.4, 12.7, 1.4, '7B1FA2', 'AGENT ORCHESTRATION — FastAPI + LangGraph + LangChain\nLangGraph Supervisor → 25+ Agents across 5 Tiers + Research + Security\nT1 Monitoring → T2 Diagnostic → T3 Risk → T4 Intervention → T5 Action\nTools: FHIR Query | Neo4j Cypher | Qdrant RAG | NL2SQL | Geospatial | Whisper Voice'),
        (0.3, 5.9, 12.7, 1.3, 'F57F17', 'DATA LAYER\nPostgreSQL 15 (FHIR R4 + PostGIS + pg_vector + multi-tenant)  |  Neo4j 5 (Clinical Knowledge Graph, 11 node types)\nQdrant (5 vector collections for RAG)  |  Redis 7 (Cache, Celery Broker, A2A PubSub)  |  Ollama/Llama 3.2 (Local LLM)  |  MinIO (Object Storage)'),
    ]
    for (l, t, w, h, c, txt) in layers:
        add_rect(slide, l, t, w, h, c, None, txt, 11, '#FFFFFF', True)

    # ============================================================
    # SLIDE 3: AI Agent System
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 0.8, '6A1B9A')
    add_textbox(slide, 0.3, 0.1, 12, 0.6, '25+ Autonomous AI Agents — 5-Tier Architecture', 28, True, '#FFFFFF')

    add_rect(slide, 0.3, 1.0, 12.7, 0.6, '4A148C', None,
             'LangGraph Supervisor — StateGraph with Conditional Routing, Parallel Execution, HITL Checkpoints',
             13, '#FFFFFF', True)

    tiers = [
        (0.3, 1.8, 2.3, 'E91E63', 'T1: MONITORING\nGlucose (LSTM CGM)\nCardiac (ECG/vitals)\nActivity (Wearable)\nTemperature (Fever)'),
        (2.8, 1.8, 2.3, 'E65100', 'T2: DIAGNOSTIC\nECG (STEMI detect)\nKidney (eGFR trend)\nImaging (AI radiology)\nLab Interpretation'),
        (5.3, 1.8, 2.7, 'F9A825', 'T3: RISK ASSESSMENT\nComorbidity Risk\nPrediction (XGBoost + RAG)\nFamily History (Neo4j)\nSDOH Risk\nML Ensemble (Attention)'),
        (8.2, 1.8, 2.3, '2E7D32', 'T4: INTERVENTION\nLifestyle Coaching\nPrescription (RAG)\nContraindication Check\nTriage & Emergency'),
        (10.7, 1.8, 2.3, '1565C0', 'T5: ACTION\nPhysician Notify\nPatient Notify\nScheduling\nEHR Integration\nBilling / RPM'),
    ]
    for (l, t, w, c, txt) in tiers:
        add_rect(slide, l, t, w, 2.0, c, None, txt, 10, '#FFFFFF', True)

    add_rect(slide, 0.3, 4.0, 6.0, 1.5, '00695C', None,
             'RESEARCH SYSTEM\nLiterature Agent (PubMed/Semantic Scholar)\nSynthesis Agent (Evidence Grading A/B/C)\nTrial Matching (ClinicalTrials.gov)\nGuideline Update Agent\nClinical Q&A (RAG)', 10, '#FFFFFF', True)
    add_rect(slide, 6.5, 4.0, 3.0, 1.5, 'B71C1C', None,
             'SECURITY\nPHI Detector (Presidio)\nPrompt Guardrails\nAudit Logger\nHallucination Detection', 10, '#FFFFFF', True)
    add_rect(slide, 9.7, 4.0, 3.3, 1.5, '4E342E', None,
             'AGENT TOOLS (LangChain)\nFHIR Query | Neo4j Cypher\nQdrant Search | NL2SQL\nNotification Dispatch\nGeospatial Hospital Routing\nWhisper Voice Transcription', 10, '#FFFFFF', True)

    add_rect(slide, 0.3, 5.7, 12.7, 0.5, '311B92', None,
             'Redis A2A Message Bus — 7 Channels | Typed Messages: ALERT, REQUEST, RESPONSE, DATA_UPDATE',
             11, '#FFFFFF', True)

    add_textbox(slide, 0.3, 6.4, 12.7, 0.8,
                'All agents traced via Langfuse  |  Celery for scheduled background jobs  |  HITL for critical clinical decisions  |  Health-literacy adapted patient messaging (5 levels)',
                12, False, '#555555', PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 4: Data Layer
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 0.8, 'E65100')
    add_textbox(slide, 0.3, 0.1, 12, 0.6, 'Data Layer — Multi-Database Architecture', 28, True, '#FFFFFF')

    dbs = [
        (0.3, 1.1, 4.0, 2.5, '1565C0',
         'PostgreSQL 15 (Port 5588)\n\n• FHIR R4 Schema (14 resources)\n• Clinical EMR Schema\n• Analytics/Population Health\n• Tenant Schema (per-org isolation)\n• HIPAA Audit Schema\n• pg_vector extension\n• PostGIS geospatial\n• Multi-tenant: schema-per-org'),
        (4.5, 1.1, 4.0, 2.5, '2E7D32',
         'Neo4j 5.12 (Bolt 7688)\n\n• 11 Node Types (Patient, Disease,\n  Medication, Gene, etc.)\n• 17 Relationship Types\n• Drug interaction detection\n• PageRank risk scoring\n• APOC + Graph Data Science\n• Tenant-scoped labels\n• Seed: 100+ drugs, diseases, guidelines'),
        (8.7, 1.1, 4.3, 2.5, '6A1B9A',
         'Qdrant Vector DB (Port 6388)\n\n5 Collections (1536-dim embeddings):\n• clinical_guidelines\n• medical_literature\n• patient_notes\n• drug_information\n• disease_knowledge\n\nRAG pipeline for agent decisions'),
    ]
    for (l, t, w, h, c, txt) in dbs:
        add_rect(slide, l, t, w, h, c, None, txt, 11, '#FFFFFF', True)

    row2 = [
        (0.3, 3.8, 3.0, 'D84315', 'Redis 7 (Port 6489)\n\nCache Layer\nCelery Broker\nA2A PubSub Bus\nSession Store\nTenant-prefixed keyspaces'),
        (3.5, 3.8, 3.2, '37474F', 'Ollama LLM (Port 11788)\n\nPrimary: Llama 3.2 (local, HIPAA)\nFallback: Claude API, GPT-4o\nEmbeddings: sentence-transformers\nAll calls traced via Langfuse\nPHI redaction before LLM'),
        (6.9, 3.8, 3.0, '00695C', 'MinIO (API 9588)\n\nMedical images\nClinical reports/PDFs\nML model artifacts\nFHIR document storage\nS3-compatible API'),
        (10.1, 3.8, 3.0, '4E342E', 'Celery + Beat\n\nDistributed task queue\nScheduled monitoring loops\nBatch analytics jobs\nReport generation\nAsync notification dispatch'),
    ]
    for (l, t, w, c, txt) in row2:
        add_rect(slide, l, t, w, 2.0, c, None, txt, 10, '#FFFFFF', True)

    add_textbox(slide, 0.3, 6.1, 12.7, 1.0,
                'All databases support tenant isolation  |  PostgreSQL: schema-per-org  |  Neo4j: tenant-scoped labels  |  Qdrant: tenant metadata filters  |  Redis: keyspace prefixes',
                13, False, '#555555', PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 5: Interoperability & Compliance
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 0.8, '00695C')
    add_textbox(slide, 0.3, 0.1, 12, 0.6, 'Healthcare Interoperability & HIPAA Compliance', 28, True, '#FFFFFF')

    items = [
        (0.3, 1.1, 4.0, 2.5, '1565C0',
         'FHIR R4 (14 Resources)\n\nPatient, Observation, Condition,\nMedicationRequest, DiagnosticReport,\nAppointment, CarePlan,\nAllergyIntolerance, Encounter,\nProcedure, Immunization,\nDocumentReference, Bundle,\nCapabilityStatement'),
        (4.5, 1.1, 4.0, 2.5, 'E65100',
         'HL7 v2 Message Processing\n\nADT — Admit/Discharge/Transfer\nORU — Lab/Observation Results\nORM — Order Messages\nMLLP Listener Endpoint\nMessage Queue + Processor\nAuto-map to FHIR resources'),
        (8.7, 1.1, 4.3, 2.5, 'B71C1C',
         'HIPAA Technical Safeguards\n\nAES-256 encryption at rest\nTLS 1.3 encryption in transit\nJWT (15min) + refresh (7 days)\nOAuth2 + TOTP 2FA\n8 RBAC roles (django-guardian)\nPHI auto-redaction (presidio)\nImmutable audit log\nRate limiting per tenant/user'),
    ]
    for (l, t, w, h, c, txt) in items:
        add_rect(slide, l, t, w, h, c, None, txt, 11, '#FFFFFF', True)

    items2 = [
        (0.3, 3.8, 3.0, '311B92', 'Coding Standards\n\nLOINC (labs/vitals)\nSNOMED CT (conditions)\nICD-10/ICD-11 (diagnoses)\nRxNorm (medications)\nCPT (billing procedures)'),
        (3.5, 3.8, 3.2, '00695C', 'MCP Protocol\n\nTool Registry (8 tools)\nContext Distribution\nPatient data aggregation\nLLM proxy with PHI filter\nAuth middleware'),
        (6.9, 3.8, 3.0, '4A148C', 'A2A Protocol\n\nAgent Card Registry\n.well-known/agent.json\nTask Delegation/Routing\nProtocol Versioning\n7-Channel Message Bus'),
        (10.1, 3.8, 3.0, '2E7D32', 'External Integrations\n\nPubMed E-utilities\nClinicalTrials.gov API\nSemantic Scholar API\nTwilio (SMS)\nWebRTC (Telemedicine)'),
    ]
    for (l, t, w, c, txt) in items2:
        add_rect(slide, l, t, w, 2.0, c, None, txt, 10, '#FFFFFF', True)

    # ============================================================
    # SLIDE 6: Monitoring & Observability
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 0.8, 'C2185B')
    add_textbox(slide, 0.3, 0.1, 12, 0.6, 'Observability & Monitoring Stack', 28, True, '#FFFFFF')

    obs = [
        (0.3, 1.1, 3.0, 2.5, 'E65100',
         'Prometheus (9390)\n\nMetrics Scraping:\n• Django request latency\n• Agent execution time\n• Celery queue depth\n• DB query time\n• LLM token usage\n• Critical alert counts'),
        (3.5, 1.1, 3.2, 2.5, '2E7D32',
         'Grafana (9391)\n\n5 Dashboards:\n• Agent Operations\n• Clinical Overview\n• System Health\n• Population Health\n• LLM Cost Tracking\n\nAuto-provisioned via YAML'),
        (6.9, 1.1, 3.0, 2.5, '6A1B9A',
         'Langfuse (3488)\n\nLLM Observability:\n• Trace all LLM calls\n• Agent evaluations\n• Cost tracking per agent\n• Faithfulness scores\n• Relevance metrics\n• Token usage analytics'),
        (10.1, 1.1, 3.0, 2.5, '1565C0',
         'OpenTelemetry\n\nDistributed Tracing:\n• Django instrumentation\n• FastAPI instrumentation\n• Psycopg2 (DB calls)\n• Redis instrumentation\n• HTTPX (external calls)\n• Custom agent_span()'),
    ]
    for (l, t, w, h, c, txt) in obs:
        add_rect(slide, l, t, w, h, c, None, txt, 11, '#FFFFFF', True)

    add_rect(slide, 0.3, 3.8, 6.0, 1.5, 'B71C1C', None,
             'Alertmanager (9393)\n\n3 Alert Rule Sets:\n• Agent alerts (failure rate, latency p99)\n• Clinical alerts (critical patient thresholds)\n• Infrastructure alerts (CPU, memory, disk)\n\nRouting: PagerDuty + Slack + Email',
             11, '#FFFFFF', True)
    add_rect(slide, 6.5, 3.8, 6.5, 1.5, '37474F', None,
             'Infrastructure — Docker Compose (20 services)\n\nAll services containerized with health checks\nNginx reverse proxy routing all traffic\nKubernetes Helm charts for production\nDocker network isolation (HIPAA)\nSecrets via environment variables\nGraceful degradation on service failure',
             11, '#FFFFFF', True)

    # ============================================================
    # SLIDE 7: Frontend & UI
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 0.8, '2E7D32')
    add_textbox(slide, 0.3, 0.1, 12, 0.6, 'Frontend — 17 Pages, Dark Mode, PWA', 28, True, '#FFFFFF')

    pages = [
        (0.3, 1.1, 3.0, 'E91E63', 'Patient Portal\n\nVitals Dashboard\nMedication Management\nAppointment Booking\nSecure Messaging\nHealth Goals (Gamified)\nEducation Resources'),
        (3.5, 1.1, 3.2, '1565C0', 'Clinician Dashboard\n\nRisk-stratified Patient List\nReal-time Vitals Monitor\nAI Recommendation Panel\nCare Gap Alerts\nHITL Decision Interface\nSmart Order Sets'),
        (6.9, 1.1, 3.0, '6A1B9A', 'Agent Control Panel\n\n25-Agent Status Grid\nExecution Traces/Logs\nAgent Activity Timeline\nLangfuse Integration\nTrigger Controls\nPerformance Metrics'),
        (10.1, 1.1, 3.0, 'E65100', 'Analytics Console\n\nPopulation Health Charts\nCohort Analysis\nRisk Score Gauges\nFairness Analysis\nWhat-If Simulator\nPredictive Models Viz'),
    ]
    for (l, t, w, c, txt) in pages:
        add_rect(slide, l, t, w, 2.3, c, None, txt, 10, '#FFFFFF', True)

    pages2 = [
        (0.3, 3.6, 3.0, '00695C', 'Research Interface\n\nNatural Language Q&A\nLiterature Search\nEvidence Synthesis\nClinical Trial Matching\nCitation Management'),
        (3.5, 3.6, 3.2, 'B71C1C', 'Admin & Settings\n\nMulti-Tenant Management\nUser Management / RBAC\nSystem Health Monitor\nBranding Configuration\nAPI Key Management'),
        (6.9, 3.6, 3.0, '37474F', 'Billing & Telemedicine\n\nRCM Dashboard\nClaims Tracking\nRPM Billing Codes\nVideo Consultation\nAI-Assisted Notes'),
        (10.1, 3.6, 3.0, '4E342E', 'Tech Stack\n\nVite Build Tool\nTailwindCSS + shadcn/ui\nZustand State Mgmt\nTanStack Query\nRecharts + D3\nWebSocket Real-time'),
    ]
    for (l, t, w, c, txt) in pages2:
        add_rect(slide, l, t, w, 2.0, c, None, txt, 10, '#FFFFFF', True)

    # ============================================================
    # SLIDE 8: Key Numbers
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, '0D47A1')
    add_textbox(slide, 1, 0.5, 11, 0.8, 'InHealth Chronic Care — By The Numbers', 32, True, '#FFFFFF', PP_ALIGN.CENTER)

    stats = [
        ('430+', 'Source Files'),
        ('20', 'Docker Services'),
        ('25+', 'AI Agents'),
        ('15', 'Django Apps'),
        ('14', 'FHIR R4 Resources'),
        ('5', 'ML Models'),
        ('5', 'Vector Collections'),
        ('5', 'Grafana Dashboards'),
        ('17', 'Frontend Pages'),
        ('8', 'Healthcare Roles'),
        ('7', 'LangChain Tools'),
        ('4', 'Notification Channels'),
    ]
    for i, (num, label) in enumerate(stats):
        row = i // 4
        col = i % 4
        l = 0.8 + col * 3.1
        t = 1.5 + row * 1.8
        add_rect(slide, l, t, 2.8, 1.5, '1565C0', '1976D2')
        add_textbox(slide, l, t + 0.15, 2.8, 0.7, num, 36, True, '#FFFFFF', PP_ALIGN.CENTER)
        add_textbox(slide, l, t + 0.85, 2.8, 0.5, label, 14, False, '#BBDEFB', PP_ALIGN.CENTER)

    add_textbox(slide, 1, 6.5, 11, 0.5,
                'Multi-Tenant  |  HIPAA-Compliant  |  FHIR R4 + HL7 v2  |  MCP + A2A Protocols  |  Kubernetes Ready',
                16, False, '#64B5F6', PP_ALIGN.CENTER)

    # Save
    output_path = 'docs/architecture/InHealth_Technical_Architecture.pptx'
    prs.save(output_path)
    print(f"PowerPoint saved to: {output_path}")
    return output_path


if __name__ == '__main__':
    create_presentation()
